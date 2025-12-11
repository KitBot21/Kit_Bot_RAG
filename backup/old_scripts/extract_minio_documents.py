#!/usr/bin/env python3
"""
MinIO에 저장된 문서 파일에서 텍스트 추출
우선순위 1+2: PDF, DOCX, TXT, XLSX, XLS, PPTX
"""

import sys
from pathlib import Path

# 프로젝트 루트를 Python 경로에 추가
sys.path.insert(0, str(Path(__file__).parent.parent))

import json
import tempfile
from typing import Dict, List
from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document
from pptx import Presentation
import openpyxl
import xlrd
from crawler.storage.minio_storage import MinIOStorage


def extract_pdf_from_bytes(file_data: bytes) -> str:
    """PDF 바이트에서 텍스트 추출"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        
        text = extract_pdf_text(tmp_path)
        Path(tmp_path).unlink()
        return text.strip()
    except Exception as e:
        return ""


def extract_docx_from_bytes(file_data: bytes) -> str:
    """DOCX 바이트에서 텍스트 추출"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        
        doc = Document(tmp_path)
        text = '\n'.join([para.text for para in doc.paragraphs])
        Path(tmp_path).unlink()
        return text.strip()
    except Exception as e:
        return ""


def extract_txt_from_bytes(file_data: bytes) -> str:
    """TXT 바이트에서 텍스트 읽기"""
    try:
        return file_data.decode('utf-8', errors='ignore').strip()
    except:
        return ""


def extract_xlsx_from_bytes(file_data: bytes) -> str:
    """XLSX 바이트에서 텍스트 추출"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        
        workbook = openpyxl.load_workbook(tmp_path, data_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"[시트: {sheet_name}]")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
                if row_text.strip():
                    text_parts.append(row_text)
        
        Path(tmp_path).unlink()
        return '\n'.join(text_parts).strip()
    except Exception as e:
        return ""


def extract_xls_from_bytes(file_data: bytes) -> str:
    """XLS 바이트에서 텍스트 추출"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.xls', delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        
        workbook = xlrd.open_workbook(tmp_path)
        text_parts = []
        
        for sheet in workbook.sheets():
            text_parts.append(f"[시트: {sheet.name}]")
            
            for row_idx in range(sheet.nrows):
                row_text = ' | '.join([str(cell.value) for cell in sheet.row(row_idx)])
                if row_text.strip():
                    text_parts.append(row_text)
        
        Path(tmp_path).unlink()
        return '\n'.join(text_parts).strip()
    except Exception as e:
        return ""


def extract_pptx_from_bytes(file_data: bytes) -> str:
    """PPTX 바이트에서 텍스트 추출"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        
        prs = Presentation(tmp_path)
        text_parts = []
        
        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"[슬라이드 {i}]")
            
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_parts.append(shape.text)
        
        Path(tmp_path).unlink()
        return '\n'.join(text_parts).strip()
    except Exception as e:
        return ""


def get_source_page_url(minio_storage, object_name: str) -> str:
    """MinIO 메타데이터에서 원본 페이지 URL 가져오기"""
    try:
        stat = minio_storage.client.stat_object("kit-attachments", object_name)
        metadata = stat.metadata
        
        # URL 디코딩
        import urllib.parse
        page_url = metadata.get('page-url', '')
        if page_url:
            page_url = urllib.parse.unquote(page_url)
        
        return page_url
    except Exception:
        return ""


def process_minio_documents(output_file: Path):
    """MinIO PDF/DOCX 파일 처리"""
    
    print("=" * 80)
    print("📦 MinIO 문서 텍스트 추출")
    print("=" * 80)
    
    # MinIO 연결
    minio = MinIOStorage(
        endpoint="localhost:9000",
        access_key="admin",
        secret_key="kitbot2025!",
        bucket_name="kit-attachments"
    )
    
    # 모든 객체 목록 가져오기
    print("\n1️⃣ MinIO 파일 목록 수집...")
    objects = list(minio.client.list_objects("kit-attachments", recursive=True))
    print(f"   총 파일: {len(objects):,}개")
    
    # PDF/DOCX/TXT/XLSX/XLS/PPTX 필터링
    target_extensions = {'.pdf', '.docx', '.txt', '.xlsx', '.xls', '.pptx'}
    target_objects = []
    
    for obj in objects:
        ext = Path(obj.object_name).suffix.lower()
        if ext in target_extensions:
            target_objects.append(obj)
    
    print(f"\n2️⃣ 처리 대상 파일:")
    ext_count = {}
    for obj in target_objects:
        ext = Path(obj.object_name).suffix.lower()
        ext_count[ext] = ext_count.get(ext, 0) + 1
    
    for ext, count in sorted(ext_count.items()):
        print(f"   {ext:10s}: {count:3d}개")
    
    print(f"\n3️⃣ 텍스트 추출 시작...")
    
    extracted_documents = []
    success_count = 0
    fail_count = 0
    
    for i, obj in enumerate(target_objects, 1):
        object_name = obj.object_name
        filename = Path(object_name).name
        ext = Path(object_name).suffix.lower()
        
        if i % 50 == 0:
            print(f"\n   진행: {i}/{len(target_objects)} ({i*100//len(target_objects)}%)")
        
        try:
            # 파일 다운로드
            response = minio.client.get_object("kit-attachments", object_name)
            file_data = response.read()
            response.close()
            response.release_conn()
            
            # 텍스트 추출
            if ext == '.pdf':
                text = extract_pdf_from_bytes(file_data)
            elif ext == '.docx':
                text = extract_docx_from_bytes(file_data)
            elif ext == '.txt':
                text = extract_txt_from_bytes(file_data)
            elif ext == '.xlsx':
                text = extract_xlsx_from_bytes(file_data)
            elif ext == '.xls':
                text = extract_xls_from_bytes(file_data)
            elif ext == '.pptx':
                text = extract_pptx_from_bytes(file_data)
            else:
                fail_count += 1
                continue
            
            if text and len(text) > 10:  # 최소 10자 이상
                # 원본 페이지 URL 가져오기
                page_url = get_source_page_url(minio, object_name)
                
                extracted_documents.append({
                    'text': text,
                    'title': filename,
                    'url': page_url if page_url else 'minio://kit-attachments/' + object_name,
                    'source_type': 'minio_document',
                    'document_name': filename,
                    'file_type': ext[1:],  # .pdf → pdf
                    'file_size': obj.size
                })
                success_count += 1
            else:
                fail_count += 1
        
        except Exception as e:
            print(f"\n   ❌ 실패: {filename} - {e}")
            fail_count += 1
            continue
    
    print(f"\n\n4️⃣ 결과:")
    print(f"   ✅ 성공: {success_count:,}개")
    print(f"   ❌ 실패: {fail_count:,}개")
    
    # CSV 저장
    if extracted_documents:
        import csv
        
        with open(output_file, 'w', encoding='utf-8', newline='') as f:
            writer = csv.DictWriter(f, fieldnames=[
                'text', 'title', 'url', 'source_type',
                'document_name', 'file_type', 'file_size'
            ])
            writer.writeheader()
            writer.writerows(extracted_documents)
        
        print(f"\n✅ 저장 완료: {output_file}")
        
        # 통계
        total_chars = sum(len(doc['text']) for doc in extracted_documents)
        avg_chars = total_chars / len(extracted_documents)
        
        print(f"\n📊 텍스트 통계:")
        print(f"   총 문서: {len(extracted_documents):,}개")
        print(f"   총 텍스트: {total_chars:,}자 ({total_chars/1000000:.1f}M자)")
        print(f"   평균 길이: {avg_chars:,.0f}자")
    else:
        print("\n⚠️  추출된 문서가 없습니다.")
    
    print("\n" + "=" * 80)
    print("🎉 완료!")
    print("=" * 80)


def main():
    output_file = Path("data/corpus_minio_documents.csv")
    process_minio_documents(output_file)


if __name__ == "__main__":
    main()
