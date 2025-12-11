#!/usr/bin/env python3
"""
MinIO 문서 추출 + 청킹 + 필터링
- 청크 크기: 1000자, 오버랩: 150자
- 불필요한 텍스트 패턴 제거
"""

import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent.parent))

import json
import tempfile
import re
import csv
from typing import Dict, List
from pdfminer.high_level import extract_text as extract_pdf_text
from docx import Document
from pptx import Presentation
import openpyxl
import xlrd
from crawler.storage.minio_storage import MinIOStorage

# 청킹 설정
CHUNK_SIZE = 1000
OVERLAP = 150
MIN_CHUNK_LENGTH = 100  # 최소 청크 길이

# 필터링 패턴
FILTER_PATTERNS = [
    r'^\s*차\s*례\s*$',
    r'^\s*목\s*차\s*$',
    r'^\s*참고문헌\s*$',
    r'^\s*부\s*록\s*$',
    r'^\s*Copyright.*$',
    r'^\s*저작권.*$',
    r'^\s*All Rights Reserved.*$',
    r'^\s*페이지\s*\d+\s*$',
    r'^\s*\d+\s*$',  # 페이지 번호만
    r'^\s*-\s*\d+\s*-\s*$',  # -1- 형식
]

def clean_text(text: str) -> str:
    """텍스트 정제 및 불필요한 패턴 제거"""
    if not text:
        return ""
    
    lines = text.split('\n')
    cleaned_lines = []
    
    for line in lines:
        # 필터링 패턴 체크
        should_skip = False
        for pattern in FILTER_PATTERNS:
            if re.match(pattern, line.strip(), re.IGNORECASE):
                should_skip = True
                break
        
        if not should_skip and line.strip():
            cleaned_lines.append(line)
    
    # 연속 공백 정리
    text = '\n'.join(cleaned_lines)
    text = re.sub(r'\s+', ' ', text)
    text = re.sub(r'\n\s*\n+', '\n', text)
    
    return text.strip()

def chunk_text(text: str) -> List[str]:
    """텍스트를 청크로 분할"""
    if len(text) <= CHUNK_SIZE:
        return [text] if text.strip() and len(text) >= MIN_CHUNK_LENGTH else []
    
    chunks = []
    start = 0
    
    while start < len(text):
        end = start + CHUNK_SIZE
        chunk = text[start:end]
        
        # 마지막 청크가 아니면 문장 경계에서 자르기
        if end < len(text):
            last_period = max(
                chunk.rfind('.'),
                chunk.rfind('!'),
                chunk.rfind('?'),
                chunk.rfind('。'),
                chunk.rfind('\n')
            )
            
            if last_period > CHUNK_SIZE * 0.5:
                end = start + last_period + 1
                chunk = text[start:end]
        
        chunk = chunk.strip()
        if chunk and len(chunk) >= MIN_CHUNK_LENGTH:
            chunks.append(chunk)
        
        start = end - OVERLAP
        if start <= 0 or start >= len(text):
            break
    
    return chunks

def extract_pdf_from_bytes(file_data: bytes) -> str:
    """PDF 바이트에서 텍스트 추출"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        
        text = extract_pdf_text(tmp_path)
        Path(tmp_path).unlink()
        return clean_text(text)
    except Exception as e:
        return ""

def extract_docx_from_bytes(file_data: bytes) -> str:
    """DOCX 바이트에서 텍스트 추출"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        
        doc = Document(tmp_path)
        text = '\n'.join([para.text for para in doc.paragraphs])
        Path(tmp_path).unlink()
        return clean_text(text)
    except Exception as e:
        return ""

def extract_txt_from_bytes(file_data: bytes) -> str:
    """TXT 바이트에서 텍스트 추출"""
    try:
        text = file_data.decode('utf-8', errors='ignore')
        return clean_text(text)
    except:
        try:
            text = file_data.decode('cp949', errors='ignore')
            return clean_text(text)
        except:
            return ""

def extract_xlsx_from_bytes(file_data: bytes) -> str:
    """XLSX 바이트에서 텍스트 추출"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        
        wb = openpyxl.load_workbook(tmp_path, data_only=True)
        texts = []
        
        for sheet in wb.worksheets:
            sheet_text = f"[시트: {sheet.title}]\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = ' | '.join([str(cell) if cell is not None else '' for cell in row])
                if row_text.strip():
                    sheet_text += row_text + '\n'
            texts.append(sheet_text)
        
        Path(tmp_path).unlink()
        return clean_text('\n'.join(texts))
    except Exception as e:
        return ""

def extract_xls_from_bytes(file_data: bytes) -> str:
    """XLS 바이트에서 텍스트 추출"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.xls', delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        
        wb = xlrd.open_workbook(tmp_path)
        texts = []
        
        for sheet in wb.sheets():
            sheet_text = f"[시트: {sheet.name}]\n"
            for row_idx in range(sheet.nrows):
                row_text = ' | '.join([str(cell.value) for cell in sheet.row(row_idx)])
                if row_text.strip():
                    sheet_text += row_text + '\n'
            texts.append(sheet_text)
        
        Path(tmp_path).unlink()
        return clean_text('\n'.join(texts))
    except Exception as e:
        return ""

def extract_pptx_from_bytes(file_data: bytes) -> str:
    """PPTX 바이트에서 텍스트 추출"""
    try:
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(file_data)
            tmp_path = tmp.name
        
        prs = Presentation(tmp_path)
        texts = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_text = f"[슬라이드 {i}]\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + '\n'
            texts.append(slide_text)
        
        Path(tmp_path).unlink()
        return clean_text('\n'.join(texts))
    except Exception as e:
        return ""

def get_source_page_url(minio_storage, filename: str) -> str:
    """MinIO 메타데이터에서 원본 URL 가져오기"""
    try:
        stat = minio_storage.client.stat_object('kit-attachments', f'attachments/{filename}')
        metadata = stat.metadata
        return metadata.get('x-amz-meta-source-url', '') or metadata.get('source-url', '')
    except:
        return ''

def process_minio_documents():
    """MinIO 문서 처리 메인 함수"""
    print("=" * 80)
    print("📤 MinIO 문서 추출 + 청킹")
    print("=" * 80)
    
    print(f"\n⚙️  설정:")
    print(f"   청크 크기: {CHUNK_SIZE}자")
    print(f"   오버랩: {OVERLAP}자")
    print(f"   최소 청크 길이: {MIN_CHUNK_LENGTH}자")
    
    # MinIO 연결
    minio_storage = MinIOStorage()
    
    # 파일 목록 가져오기
    print(f"\n📂 파일 목록 가져오는 중...")
    objects = minio_storage.client.list_objects('kit-attachments', prefix='attachments/', recursive=True)
    
    # 처리할 파일 필터링
    target_extensions = {'.pdf', '.docx', '.txt', '.xlsx', '.xls', '.pptx'}
    files_to_process = []
    
    for obj in objects:
        filename = obj.object_name.replace('attachments/', '')
        ext = Path(filename).suffix.lower()
        if ext in target_extensions:
            files_to_process.append((filename, ext, obj.size))
    
    print(f"   총 파일: {len(files_to_process)}개")
    
    # 파일 타입별 통계
    from collections import Counter
    ext_counts = Counter([ext for _, ext, _ in files_to_process])
    for ext, count in sorted(ext_counts.items()):
        print(f"   {ext}: {count}개")
    
    # 처리
    results = []
    success_count = 0
    failed_count = 0
    total_chunks = 0
    
    print(f"\n⏳ 처리 중...")
    
    for i, (filename, ext, size) in enumerate(files_to_process, 1):
        if i % 50 == 0:
            print(f"   진행: {i}/{len(files_to_process)} ({i/len(files_to_process)*100:.0f}%)")
        
        try:
            # 파일 다운로드
            file_data = minio_storage.download_file(filename)
            
            # 텍스트 추출
            text = ""
            if ext == '.pdf':
                text = extract_pdf_from_bytes(file_data)
            elif ext == '.docx':
                text = extract_docx_from_bytes(file_data)
            elif ext == '.txt':
                text = extract_txt_from_bytes(file_data)
            elif ext == '.xlsx':
                text = extract_xlsx_from_bytes(file_data)
            elif ext == '.xls':
                text = extract_xls_from_bytes(file_data)
            elif ext == '.pptx':
                text = extract_pptx_from_bytes(file_data)
            
            if not text or len(text) < MIN_CHUNK_LENGTH:
                failed_count += 1
                continue
            
            # 청킹
            chunks = chunk_text(text)
            
            if not chunks:
                failed_count += 1
                continue
            
            # 원본 URL 가져오기
            source_url = get_source_page_url(minio_storage, filename)
            
            # 각 청크를 개별 레코드로 저장
            for chunk_idx, chunk in enumerate(chunks):
                results.append({
                    'text': chunk,
                    'title': filename,
                    'url': source_url,
                    'source_type': 'minio_document',
                    'document_name': f"{filename}_chunk{chunk_idx}",
                    'file_type': ext,
                    'file_size': size,
                    'chunk_index': chunk_idx,
                    'total_chunks': len(chunks)
                })
            
            success_count += 1
            total_chunks += len(chunks)
            
        except Exception as e:
            failed_count += 1
    
    print(f"\n4️⃣ 결과:")
    print(f"   ✅ 성공: {success_count}개 파일")
    print(f"   ❌ 실패: {failed_count}개 파일")
    print(f"   📊 총 청크: {total_chunks}개")
    
    # CSV 저장
    output_path = Path("data/corpus_minio_documents.csv")
    
    if results:
        with output_path.open('w', encoding='utf-8', newline='') as f:
            fieldnames = ['text', 'title', 'url', 'source_type', 'document_name', 
                         'file_type', 'file_size', 'chunk_index', 'total_chunks']
            writer = csv.DictWriter(f, fieldnames=fieldnames)
            writer.writeheader()
            writer.writerows(results)
        
        print(f"\n✅ 저장 완료: {output_path}")
        
        # 통계
        total_text_length = sum(len(r['text']) for r in results)
        avg_chunk_length = total_text_length / len(results) if results else 0
        
        print(f"\n📊 텍스트 통계:")
        print(f"   총 문서: {len(results)}개")
        print(f"   총 텍스트: {total_text_length:,}자")
        print(f"   평균 청크 길이: {avg_chunk_length:.0f}자")
    
    print("\n" + "=" * 80)
    print("🎉 완료!")
    print("=" * 80)

if __name__ == "__main__":
    process_minio_documents()
