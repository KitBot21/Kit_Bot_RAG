#!/usr/bin/env python3
"""
첨부파일 처리: PDF, Word, Excel, PPT 등을 텍스트로 변환하여 corpus에 추가

지원 방식:
1. 로컬 파일 (기본): data/attachments/ 디렉토리에 파일 직접 저장
2. MinIO/S3 (옵션): Object Storage에서 파일 다운로드 후 처리
"""
import csv
import hashlib
import re
import argparse
import os
from pathlib import Path
from datetime import datetime
from dotenv import load_dotenv

# .env 파일 로드
load_dotenv()

# PDF 처리
try:
    import PyPDF2
    HAS_PDF = True
except ImportError:
    HAS_PDF = False
    print("⚠️  PyPDF2 미설치: pip install PyPDF2")

# Word/Excel/PPT 처리
try:
    import docx
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False
    print("⚠️  python-docx 미설치: pip install python-docx")

try:
    import openpyxl
    HAS_EXCEL = True
except ImportError:
    HAS_EXCEL = False
    print("⚠️  openpyxl 미설치: pip install openpyxl")

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("⚠️  python-pptx 미설치: pip install python-pptx")

PROJECT_ROOT = Path(__file__).resolve().parents[1]
ATTACHMENTS_DIR = PROJECT_ROOT / "data" / "attachments"
OUT_CSV = PROJECT_ROOT / "data" / "corpus_attachments.csv"

# MinIO/S3 지원 (옵션)
try:
    from minio import Minio
    HAS_MINIO = True
except ImportError:
    HAS_MINIO = False

# 청킹 설정
CHARS = 800
OVERLAP = 100

def sha256(s: str) -> str:
    return hashlib.sha256(s.encode("utf-8")).hexdigest()

def clean(s: str) -> str:
    """텍스트 정제"""
    s = s.replace("\u00a0", " ").replace("\t", " ")
    s = re.sub(r"[ \t]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()

def chunk_text(full_text: str, size=CHARS, overlap=OVERLAP):
    """텍스트를 청크로 분할"""
    full_text = clean(full_text)
    n = len(full_text)
    if n == 0:
        return []
    
    out = []
    i = 0
    while i < n:
        j = min(n, i + size)
        k = j
        
        # 문장 경계 보정
        sentence_endings = ".!?。\n"
        for off in range(200):
            if j + off < n and full_text[j + off:j + off + 1] in sentence_endings:
                k = j + off + 1
                break
        
        chunk_text = full_text[i:k]
        chunk_len = len(chunk_text)
        
        # 너무 짧은 청크는 건너뛰기
        if chunk_len < 100:
            if k >= n and chunk_len >= 40:
                out.append((chunk_text, i, k))
            else:
                i = k
                continue
        else:
            out.append((chunk_text, i, k))
        
        if k >= n:
            break
        i = max(0, k - overlap)
    
    return out

def guess_lang(s: str) -> str:
    return "ko" if re.search(r"[가-힣]", s) else "en"

def extract_pdf_text(file_path: Path) -> str:
    """PDF에서 텍스트 추출"""
    if not HAS_PDF:
        return ""
    
    try:
        text = []
        with open(file_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
        return "\n\n".join(text)
    except Exception as e:
        print(f"  ⚠️  PDF 처리 실패 ({file_path.name}): {e}")
        return ""

def extract_docx_text(file_path: Path) -> str:
    """Word 문서에서 텍스트 추출"""
    if not HAS_DOCX:
        return ""
    
    try:
        doc = docx.Document(file_path)
        text = []
        for para in doc.paragraphs:
            if para.text.strip():
                text.append(para.text)
        
        # 표 내용도 추출
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    text.append(row_text)
        
        return "\n\n".join(text)
    except Exception as e:
        print(f"  ⚠️  DOCX 처리 실패 ({file_path.name}): {e}")
        return ""

def extract_excel_text(file_path: Path) -> str:
    """Excel에서 텍스트 추출"""
    if not HAS_EXCEL:
        return ""
    
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        text = []
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text.append(f"[{sheet_name}]")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                if row_text.strip():
                    text.append(row_text)
        
        return "\n\n".join(text)
    except Exception as e:
        print(f"  ⚠️  Excel 처리 실패 ({file_path.name}): {e}")
        return ""

def extract_pptx_text(file_path: Path) -> str:
    """PowerPoint에서 텍스트 추출"""
    if not HAS_PPTX:
        return ""
    
    try:
        prs = Presentation(file_path)
        text = []
        
        for i, slide in enumerate(prs.slides, 1):
            text.append(f"[슬라이드 {i}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        
        return "\n\n".join(text)
    except Exception as e:
        print(f"  ⚠️  PPTX 처리 실패 ({file_path.name}): {e}")
        return ""

def extract_text_from_file(file_path: Path) -> str:
    """파일 형식에 따라 텍스트 추출"""
    ext = file_path.suffix.lower()
    
    if ext == '.pdf':
        return extract_pdf_text(file_path)
    elif ext in ['.docx', '.doc']:
        return extract_docx_text(file_path)
    elif ext in ['.xlsx', '.xls']:
        return extract_excel_text(file_path)
    elif ext in ['.pptx', '.ppt']:
        return extract_pptx_text(file_path)
    elif ext == '.txt':
        try:
            return file_path.read_text(encoding='utf-8', errors='ignore')
        except:
            return ""
    else:
        print(f"  ⚠️  지원하지 않는 파일 형식: {ext}")
        return ""

def download_from_minio(minio_config: dict, target_dir: Path):
    """
    MinIO에서 파일 다운로드
    
    Args:
        minio_config: MinIO 설정 {'endpoint', 'access_key', 'secret_key', 'bucket'}
        target_dir: 다운로드 대상 디렉토리
    """
    if not HAS_MINIO:
        print("❌ MinIO 라이브러리 미설치: pip install minio")
        return False
    
    try:
        print(f"\n🔗 MinIO 연결 중...")
        print(f"   Endpoint: {minio_config['endpoint']}")
        print(f"   Bucket: {minio_config['bucket']}")
        
        client = Minio(
            minio_config['endpoint'],
            access_key=minio_config['access_key'],
            secret_key=minio_config['secret_key'],
            secure=minio_config.get('secure', True)
        )
        
        # 버킷 존재 확인
        bucket = minio_config['bucket']
        if not client.bucket_exists(bucket):
            print(f"❌ 버킷이 존재하지 않습니다: {bucket}")
            return False
        
        # 파일 목록 가져오기
        objects = client.list_objects(bucket, recursive=True)
        
        target_dir.mkdir(parents=True, exist_ok=True)
        downloaded = 0
        
        for obj in objects:
            # 지원하는 파일 형식만 다운로드
            ext = Path(obj.object_name).suffix.lower()
            if ext in {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.txt'}:
                local_path = target_dir / Path(obj.object_name).name
                
                print(f"  📥 다운로드: {obj.object_name}")
                client.fget_object(bucket, obj.object_name, str(local_path))
                downloaded += 1
        
        print(f"\n✅ MinIO에서 {downloaded}개 파일 다운로드 완료")
        return True
        
    except Exception as e:
        print(f"❌ MinIO 연결 실패: {e}")
        return False

def process_attachments():
    """첨부파일 처리 메인 함수"""
    print("=" * 80)
    print("📎 첨부파일 처리 시작")
    print("=" * 80)
    
    # 명령행 인자 파싱
    parser = argparse.ArgumentParser(description='첨부파일 처리 및 corpus 생성')
    parser.add_argument('--source', choices=['local', 'minio'], default='local',
                        help='파일 소스 (local: 로컬 디렉토리, minio: MinIO/S3)')
    parser.add_argument('--minio-endpoint', help='MinIO endpoint (기본값: .env의 MINIO_ENDPOINT)')
    parser.add_argument('--minio-access-key', help='MinIO access key (기본값: .env의 MINIO_ACCESS_KEY)')
    parser.add_argument('--minio-secret-key', help='MinIO secret key (기본값: .env의 MINIO_SECRET_KEY)')
    parser.add_argument('--minio-bucket', help='MinIO bucket 이름 (기본값: .env의 MINIO_BUCKET)')
    parser.add_argument('--minio-secure', action='store_true', help='HTTPS 사용 (기본값: .env의 MINIO_SECURE)')
    args = parser.parse_args()
    
    # MinIO에서 다운로드
    if args.source == 'minio':
        # .env 파일 또는 명령행 인자에서 설정 읽기
        minio_config = {
            'endpoint': args.minio_endpoint or os.getenv('MINIO_ENDPOINT'),
            'access_key': args.minio_access_key or os.getenv('MINIO_ACCESS_KEY'),
            'secret_key': args.minio_secret_key or os.getenv('MINIO_SECRET_KEY'),
            'bucket': args.minio_bucket or os.getenv('MINIO_BUCKET', 'kit-attachments'),
            'secure': args.minio_secure or os.getenv('MINIO_SECURE', 'false').lower() == 'true'
        }
        
        if not all([minio_config['endpoint'], minio_config['access_key'], 
                    minio_config['secret_key']]):
            print("❌ MinIO 설정이 부족합니다.")
            print("\n.env 파일에 다음을 추가하거나 명령행 옵션을 사용하세요:")
            print("  MINIO_ENDPOINT=localhost:9000")
            print("  MINIO_ACCESS_KEY=your_access_key")
            print("  MINIO_SECRET_KEY=your_secret_key")
            print("  MINIO_BUCKET=kit-attachments")
            print("  MINIO_SECURE=false")
            return
        
        if not download_from_minio(minio_config, ATTACHMENTS_DIR):
            return
    
    # 로컬 디렉토리 확인
    if not ATTACHMENTS_DIR.exists():
        print(f"\n❌ 첨부파일 디렉토리가 없습니다: {ATTACHMENTS_DIR}")
        print(f"\n💡 사용 방법:")
        print(f"   1. 로컬 파일 사용:")
        print(f"      mkdir -p {ATTACHMENTS_DIR}")
        print(f"      # 첨부파일들을 {ATTACHMENTS_DIR}에 복사")
        print(f"      python3 scripts/process_attachments.py")
        print(f"\n   2. MinIO 사용:")
        print(f"      python3 scripts/process_attachments.py --source minio \\")
        print(f"        --minio-endpoint localhost:9000 \\")
        print(f"        --minio-access-key YOUR_KEY \\")
        print(f"        --minio-secret-key YOUR_SECRET \\")
        print(f"        --minio-bucket kit-attachments")
        ATTACHMENTS_DIR.mkdir(parents=True, exist_ok=True)
        print(f"\n   📁 디렉토리 생성 완료: {ATTACHMENTS_DIR}")
        return
    
    # 지원 파일 형식
    supported_extensions = {'.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.txt'}
    files = [f for f in ATTACHMENTS_DIR.rglob("*") if f.is_file() and f.suffix.lower() in supported_extensions]
    
    if not files:
        print(f"\n⚠️  처리할 파일이 없습니다.")
        print(f"   지원 형식: {', '.join(supported_extensions)}")
        return
    
    print(f"\n📊 발견된 파일: {len(files)}개")
    
    rows = []
    processed = 0
    skipped = 0
    
    for file_path in sorted(files):
        print(f"\n📄 처리 중: {file_path.name}")
        
        # 텍스트 추출
        text = extract_text_from_file(file_path)
        
        if not text or len(text) < 40:
            print(f"  ⚠️  텍스트가 너무 짧음 (길이: {len(text)})")
            skipped += 1
            continue
        
        # 문서 정보
        doc_id = re.sub(r"\W+", "_", file_path.stem.lower()).strip("_")
        title = file_path.stem
        file_type = file_path.suffix.lower().lstrip('.')
        
        # 메타데이터
        stat = file_path.stat()
        modified_date = datetime.fromtimestamp(stat.st_mtime).strftime("%Y-%m-%d")
        
        # 청킹
        chunks = chunk_text(text)
        
        if not chunks:
            print(f"  ⚠️  청크 생성 실패")
            skipped += 1
            continue
        
        print(f"  ✅ {len(chunks)}개 청크 생성 (총 {len(text):,}자)")
        
        # Corpus에 추가
        for idx, (chunk, s0, s1) in enumerate(chunks, start=1):
            if len(chunk) < 40 or len(chunk) > 5000:
                continue
            
            rows.append({
                "chunk_id": f"{doc_id}_{idx:04d}",
                "doc_id": doc_id,
                "text": chunk,
                "title": title,
                "url": "",  # 첨부파일은 URL 없음
                "canonical_url": "",
                "snapshot_url": str(file_path.relative_to(PROJECT_ROOT)),
                "domain": "attachment",
                "source_path": f"/attachments/{file_path.name}",
                "section": "attachments",
                "accessed_at": modified_date,
                "lastmod": modified_date,
                "publisher": "금오공과대학교",
                "selector": file_type,
                "char_start": s0,
                "char_end": s1,
                "chunk_sha256": sha256(chunk),
                "page_sha256": sha256(text),
                "lang": guess_lang(chunk),
                "tags": f"attachment,{file_type}"
            })
        
        processed += 1
    
    # CSV 저장
    if rows:
        with OUT_CSV.open("w", encoding="utf-8", newline="") as f:
            writer = csv.DictWriter(f, fieldnames=[
                "chunk_id", "doc_id", "text", "title", "url", "canonical_url", "snapshot_url",
                "domain", "source_path", "section", "accessed_at", "lastmod", "publisher",
                "selector", "char_start", "char_end", "chunk_sha256", "page_sha256", "lang", "tags"
            ])
            writer.writeheader()
            writer.writerows(rows)
        
        print("\n" + "=" * 80)
        print("✅ 첨부파일 처리 완료!")
        print("=" * 80)
        print(f"  처리된 파일: {processed}개")
        print(f"  건너뛴 파일: {skipped}개")
        print(f"  총 청크 수: {len(rows)}개")
        print(f"  저장 위치: {OUT_CSV}")
        print("\n💡 다음 단계:")
        print("  1. corpus_filtered.csv와 corpus_attachments.csv 병합")
        print("  2. python3 scripts/regenerate_embeddings.py 실행")
        print("  3. python3 scripts/ingest_multi.py 실행")
    else:
        print("\n⚠️  생성된 청크가 없습니다.")

def main():
    process_attachments()

if __name__ == "__main__":
    main()
